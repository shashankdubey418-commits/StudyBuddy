"""
parser.py
Extracts raw text from PDF, DOCX, and PPTX files.
"""

import os
import pdfplumber
from docx import Document
from pptx import Presentation


def extract_text(file_path: str) -> str:
    """Detects file type from extension and routes to the right extractor."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Use PDF, DOCX, or PPTX.")


def _extract_pdf(file_path: str) -> str:
    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"[Page {i + 1}]\n{page_text}")
    return "\n\n".join(text_parts)


def _extract_docx(file_path: str) -> str:
    doc = Document(file_path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]

    # Include table content too, since study notes often have tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)

    return "\n".join(parts)


def _extract_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_parts = []

    for i, slide in enumerate(prs.slides):
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        slide_lines.append(line)
        if slide_lines:
            text_parts.append(f"[Slide {i + 1}]\n" + "\n".join(slide_lines))

    return "\n\n".join(text_parts)
